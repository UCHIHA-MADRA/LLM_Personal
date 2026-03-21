"""
Knowledge Base — RAG system for document Q&A.
Ingest your documents, ask questions, get answers grounded in your files.
Uses ChromaDB + sentence-transformers. Everything stored locally.
"""

import os
import logging
import builtins
from pathlib import Path
from typing import Optional, List, Dict, Any

from . import config  # type: ignore

logger = logging.getLogger(__name__)

# Lazy imports for optional dependencies
_chromadb = None
_sentence_transformer = None


def _get_chromadb():
    global _chromadb
    if _chromadb is None:
        try:
            import chromadb  # type: ignore
            _chromadb = chromadb
        except ImportError:
            raise ImportError(
                "chromadb is not installed. Install it: pip install chromadb"
            )
    return _chromadb


def _get_embedding_model():
    global _sentence_transformer
    if _sentence_transformer is None:
        try:
            from sentence_transformers import SentenceTransformer  # type: ignore
            print(f"📦 Loading embedding model: {config.EMBEDDING_MODEL_NAME}...")
            _sentence_transformer = SentenceTransformer(config.EMBEDDING_MODEL_NAME)
            print("✅ Embedding model loaded.")
        except ImportError:
            raise ImportError(
                "sentence-transformers is not installed. Install it: pip install sentence-transformers"
            )
    return _sentence_transformer


class EmbeddingFunction:
    """ChromaDB-compatible embedding function using sentence-transformers."""

    name: str = "sentence-transformers"  # Required by ChromaDB >= 0.4

    def __init__(self):
        self.model = _get_embedding_model()

    def __call__(self, input: List[str]) -> List[List[float]]:
        embeddings = self.model.encode(input, show_progress_bar=False)
        return embeddings.tolist()


def _chunk_text(text: str, chunk_size: Optional[int] = None, overlap: Optional[int] = None) -> List[str]:
    """Split text into overlapping chunks."""
    chunk_size = chunk_size or config.CHUNK_SIZE
    overlap = overlap or config.CHUNK_OVERLAP
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size  # type: ignore[operator]
        chunk = text[start:end]  # type: ignore[index]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap  # type: ignore[operator]
    return chunks


class KnowledgeBase:
    """
    RAG system for document Q&A.
    Ingest documents → they're chunked, embedded, stored in ChromaDB.
    Query → relevant chunks retrieved → injected into LLM context.
    """

    def __init__(self, db_dir: Optional[str] = None):
        self.db_dir = Path(db_dir) if db_dir else config.KNOWLEDGE_DB_DIR
        self.db_dir.mkdir(parents=True, exist_ok=True)
        self._client: Any = None
        self._collection: Any = None
        self._embedding_fn: Any = None
        self._initialized = False

    def _ensure_initialized(self):
        """Lazy initialization of ChromaDB and embeddings."""
        if self._initialized:
            return

        chromadb = _get_chromadb()
        self._embedding_fn = EmbeddingFunction()

        self._client = chromadb.PersistentClient(path=str(self.db_dir))
        self._collection = self._client.get_or_create_collection(
            name="personal_knowledge",
            embedding_function=self._embedding_fn,
            metadata={"hnsw:space": "cosine"},
        )
        self._initialized = True
        logger.info(f"Knowledge base initialized at {self.db_dir}")

    def add_text(self, text: str, source: str = "manual", metadata: Optional[Dict] = None):
        """
        Add raw text to the knowledge base.

        Args:
            text: The text content to add
            source: Source identifier (filename, URL, etc.)
            metadata: Optional metadata dict
        """
        self._ensure_initialized()
        chunks = _chunk_text(text)

        if not chunks:
            print("⚠️ No content to add.")
            return 0

        ids = [f"{source}_{i}" for i in range(len(chunks))]
        metadatas = [{"source": source, "chunk_index": i, **(metadata or {})} for i in range(len(chunks))]

        self._collection.upsert(
            documents=chunks,
            ids=ids,
            metadatas=metadatas,
        )

        print(f"✅ Added {len(chunks)} chunks from '{source}'")
        return len(chunks)

    def add_file(self, file_path: str) -> int:
        """
        Add a file to the knowledge base.
        Supports: .txt, .md, .py, .json, .csv, .pdf, .docx, .xlsx, .pptx
        """
        path = Path(file_path)
        if not path.exists():
            print(f"❌ File not found: {file_path}")
            return 0

        ext = path.suffix.lower()
        text = ""

        if ext == ".pdf":
            text = self._read_pdf(path)
        elif ext in (".docx", ".doc"):
            text = self._read_docx(path)
        elif ext in (".xlsx", ".xls"):
            text = self._read_xlsx(path)
        elif ext in (".pptx", ".ppt"):
            text = self._read_pptx(path)
        elif ext in (".txt", ".md", ".py", ".js", ".ts", ".java", ".cpp",
                      ".c", ".h", ".go", ".rs", ".rb", ".php", ".css",
                      ".html", ".xml", ".yaml", ".yml", ".toml", ".ini",
                      ".cfg", ".conf", ".sh", ".bat", ".ps1", ".sql"):
            text = path.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".json":
            text = path.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".csv":
            text = path.read_text(encoding="utf-8", errors="ignore")
        else:
            # Try reading as text
            try:
                text = path.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                print(f"❌ Unsupported file type: {ext}")
                return 0

        if not text.strip():
            print(f"⚠️ File is empty: {file_path}")
            return 0

        return self.add_text(text, source=path.name)

    def _read_pdf(self, path: Path) -> str:
        """Extract text from a PDF file."""
        try:
            from PyPDF2 import PdfReader  # type: ignore
            reader = PdfReader(str(path))
            pages = []
            for i, page in enumerate(reader.pages):
                try:
                    text = page.extract_text()
                    if text:
                        pages.append(text)
                except Exception as page_err:
                    logger.warning(f"Failed to extract text from page {i} of {path.name}: {page_err}")
            return "\n\n".join(pages)
        except ImportError:
            print("⚠️ PyPDF2 not installed. Install it: pip install PyPDF2")
            return ""
        except Exception as e:
            logger.exception(f"PDF extraction failed for {path.name}: {e}")
            return ""

    def _read_docx(self, path: Path) -> str:
        """Extract text from a DOCX file."""
        try:
            from docx import Document  # type: ignore
            doc = Document(str(path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except ImportError:
            print("⚠️ python-docx not installed. Install it: pip install python-docx")
            return ""
        except Exception as e:
            print(f"❌ Failed to read DOCX: {e}")
            return ""

    def _read_xlsx(self, path: Path) -> str:
        """Extract text from an XLSX file."""
        try:
            from openpyxl import load_workbook  # type: ignore
            wb = load_workbook(str(path), read_only=True, data_only=True)
            rows = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows.append(f"[Sheet: {sheet}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
            wb.close()
            return "\n".join(rows)
        except ImportError:
            print("⚠️ openpyxl not installed. Install it: pip install openpyxl")
            return ""
        except Exception as e:
            print(f"❌ Failed to read XLSX: {e}")
            return ""

    def _read_pptx(self, path: Path) -> str:
        """Extract text from a PPTX file."""
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(str(path))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():  # type: ignore[attr-defined]
                        texts.append(shape.text)  # type: ignore[attr-defined]
                if texts:
                    slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
            return "\n\n".join(slides_text)
        except ImportError:
            print("⚠️ python-pptx not installed. Install it: pip install python-pptx")
            return ""
        except Exception as e:
            print(f"❌ Failed to read PPTX: {e}")
            return ""

    def query(self, question: str, n_results: Optional[int] = None) -> str:
        """
        Query the knowledge base and return relevant context.

        Args:
            question: The user's question
            n_results: Number of relevant chunks to retrieve

        Returns:
            Formatted context string with relevant document chunks
        """
        self._ensure_initialized()
        n_results = n_results or config.TOP_K_RESULTS

        if self._collection.count() == 0:
            return ""

        results = self._collection.query(
            query_texts=[question],
            n_results=min(n_results, self._collection.count()),
        )

        if not results["documents"] or not results["documents"][0]:
            return ""

        # Format context
        context_parts = []
        for i, (doc, meta) in enumerate(
            zip(results["documents"][0], results["metadatas"][0])
        ):
            source = meta.get("source", "unknown")
            context_parts.append(f"[Source: {source}]\n{doc}")

        return "\n\n---\n\n".join(context_parts)

    def get_stats(self) -> Dict[str, Any]:
        """Get statistics about the knowledge base."""
        self._ensure_initialized()
        count = self._collection.count()
        return {
            "total_chunks": count,
            "db_path": str(self.db_dir),
        }

    def clear(self):
        """Clear all documents from the knowledge base."""
        self._ensure_initialized()
        # Delete and recreate collection
        self._client.delete_collection("personal_knowledge")
        self._collection = self._client.get_or_create_collection(
            name="personal_knowledge",
            embedding_function=self._embedding_fn,
            metadata={"hnsw:space": "cosine"},
        )
        print("🗑️ Knowledge base cleared.")

    def delete_source(self, source_name: str) -> int:
        """
        Delete all chunks belonging to a specific document source.

        Args:
            source_name: The source identifier (filename) to delete

        Returns:
            Number of chunks deleted
        """
        self._ensure_initialized()
        # Find all chunk IDs with this source
        all_data = self._collection.get(where={"source": source_name})
        ids_to_delete = all_data.get("ids", [])

        if not ids_to_delete:
            print(f"⚠️ No chunks found for source '{source_name}'")
            return 0

        self._collection.delete(ids=ids_to_delete)
        print(f"🗑️ Deleted {len(ids_to_delete)} chunks from '{source_name}'")
        return len(ids_to_delete)

    def list_sources(self) -> List[str]:
        """List all unique document sources in the knowledge base."""
        self._ensure_initialized()
        if self._collection.count() == 0:
            return []

        # Get all metadata
        all_data = self._collection.get()
        sources = set()
        for meta in all_data.get("metadatas", []):
            if meta and "source" in meta:
                sources.add(meta["source"])
        return sorted(sources)
